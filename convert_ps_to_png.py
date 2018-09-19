#!/usr/bin/env python
from __future__ import division
from pptx import Presentation
from pptx.util import Inches
import glob,os,sys,multiprocessing

nprocs = 5
cmd = 'convert -density 100 {input} -flatten -rotate 90 {output}'
def convert(image):
   global cmd
   output = image.replace('.ps','.png')
   outputs.append(output)
   os.system(cmd.format(input=image,output=output))

if len(sys.argv) == 2:
   images = sorted(glob.glob(sys.argv[1]))
else:
   images = sorted(glob.glob('./*.ps'))

print('processing %s images' % len(images))
outputs = []
p = multiprocessing.Pool()

percent = int(len(images)*0.025)
for i, _ in enumerate(p.imap_unordered(convert, images), 1):
   if i % percent == 0:
      sys.stderr.write('\rdone {0:5.2f}%'.format((i/len(images))*100))

'''
for i in range(len(images)):
   image = images[i]
   if i % five_percent == 0:
      print(' %4d%% done' % int(i*1./len(images)*100.))
   output = image.replace('.ps','.png')
   outputs.append(output)
   os.system(cmd.format(input=image,output=output))
'''

print('making presentation')

outputs = sorted(glob.glob('*.png'))

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

slide_width = prs.slide_width
slide_height = prs.slide_height

for i in range(0,len(outputs),4):
   # add new slide
   slide = prs.slides.add_slide(blank_slide_layout)

   # top-left plot
   if i < len(outputs):
      top = 0.01 * slide_height
      left = 0.01 * slide_width
      height = 0.49 * slide_height
      width = 0.49 * slide_width
      slide.shapes.add_picture(outputs[i],left,top,width,height)

   # top-right plot
   if i + 1 < len(outputs):
      top = 0.01 * slide_height
      left = 0.5 * slide_width
      height = 0.49 * slide_height
      width = 0.49 * slide_width
      slide.shapes.add_picture(outputs[i + 1],left,top,width,height)

   # bottom-left plot
   if i + 2 < len(outputs):
      top = 0.5 * slide_height
      left = 0.01 * slide_width
      height = 0.49 * slide_height
      width = 0.49 * slide_width
      slide.shapes.add_picture(outputs[i + 2],left,top,width,height)

   # bottom-right plot
   if i + 3 < len(outputs):
      top = 0.5 * slide_height
      left = 0.5 * slide_width
      height = 0.49 * slide_height
      width = 0.49 * slide_width
      slide.shapes.add_picture(outputs[i + 3],left,top,width,height)


prs.save('test.pptx')


